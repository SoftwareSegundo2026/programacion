# -*- coding: utf-8 -*-
"""
Genera la presentación PPTX del taller:
"Consumo de APIs en Next.js — TanStack Query · React Hook Form · Zod"

Uso:
    python3 generar_ppt.py          # crea tanstack-query-rhf.pptx
    # convertir a PDF:
    soffice --headless --convert-to pdf --outdir . tanstack-query-rhf.pptx
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------- paleta
DARK = RGBColor(0x0F, 0x17, 0x2A)      # slate-900
DARK_SOFT = RGBColor(0x1E, 0x29, 0x3B)
ACCENT = RGBColor(0x63, 0x66, 0xF1)    # indigo-500
ACCENT_SOFT = RGBColor(0xE0, 0xE7, 0xFF)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0xF1, 0xF5, 0xF9)
LINE = RGBColor(0xE2, 0xE8, 0xF0)

FONT_TITLE = "Segoe UI"
FONT_BODY = "Segoe UI"
FONT_MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _slide():
    return prs.slides.add_slide(BLANK)


def _rect(slide, left, top, width, height, color, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, left, top, width, height)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _text(slide, left, top, width, height, runs, size=18, color=DARK_SOFT,
          bold=False, align=PP_ALIGN.LEFT, font=FONT_BODY, space_after=6,
          line_spacing=1.0):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    if isinstance(runs, str):
        runs = [(runs, {})]
    for text, opts in runs:
        r = p.add_run()
        r.text = text
        r.font.name = opts.get("font", font)
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", bold)
        r.font.italic = opts.get("italic", False)
        r.font.color.rgb = opts.get("color", color)
    return tb


def _bullets(slide, items, left=Inches(0.9), top=Inches(1.9),
             width=Inches(11.5), height=Inches(4.9)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        kind = item[0]
        if kind == "bullet":
            txt = item[1]
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(10)
            p.line_spacing = 1.15
            r = p.add_run()
            r.text = "▪  "
            r.font.name = FONT_BODY
            r.font.size = Pt(18)
            r.font.color.rgb = ACCENT
            r.font.bold = True
            if isinstance(txt, str):
                txt = [(txt, {})]
            for text, opts in txt:
                r2 = p.add_run()
                r2.text = text
                r2.font.name = FONT_BODY
                r2.font.size = Pt(opts.get("size", 18))
                r2.font.bold = opts.get("bold", False)
                r2.font.color.rgb = opts.get("color", DARK_SOFT)
        elif kind == "code":
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(8)
            p.line_spacing = 1.1
            r = p.add_run()
            r.text = "    " + item[1].replace("\n", "\n    ")
            r.font.name = FONT_MONO
            r.font.size = Pt(14)
            r.font.color.rgb = DARK
        elif kind == "spacer":
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(6)
        elif kind == "note":
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(6)
            r = p.add_run()
            r.text = item[1]
            r.font.name = FONT_BODY
            r.font.size = Pt(14)
            r.font.italic = True
            r.font.color.rgb = MUTED
    return tb


def _content_slide(title, subtitle=None, items=None):
    s = _slide()
    _rect(s, 0, 0, SW, SH, WHITE)
    _text(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.9),
          [(title, {"bold": True, "size": 30, "color": DARK})])
    _rect(s, Inches(0.9), Inches(1.35), Inches(1.1), Inches(0.07), ACCENT)
    if subtitle:
        _text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.4),
              [(subtitle, {"size": 15, "color": MUTED})])
    if items:
        _bullets(s, items)
    _text(s, Inches(0.9), Inches(7.05), Inches(11.5), Inches(0.3),
          "Taller Frontend · TanStack Query + React Hook Form + Zod",
          size=11, color=MUTED)
    return s


def _code_box(slide, text, top, left=Inches(0.9), width=Inches(11.5)):
    lines = text.split("\n")
    h = Inches(0.34 * len(lines) + 0.3)
    box = _rect(slide, left, top, width, h, CODE_BG, MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.12)
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = ln
        r.font.name = FONT_MONO
        r.font.size = Pt(14)
        r.font.color.rgb = DARK
    return box


# ------------------------------------------------------------ slide 1 · título
s = _slide()
_rect(s, 0, 0, SW, SH, DARK)
_rect(s, Inches(0.9), Inches(2.35), Inches(1.6), Inches(0.09), ACCENT)
_text(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.6),
      [("Consumo de APIs en Next.js", {"bold": True, "size": 40, "color": WHITE})])
_text(s, Inches(0.9), Inches(3.35), Inches(11.5), Inches(0.8),
      [("TanStack Query · React Hook Form · Zod", {"bold": True, "size": 26, "color": ACCENT_SOFT})])
_text(s, Inches(0.9), Inches(4.35), Inches(11.5), Inches(0.9),
      [("Por qué se aplican: datos con caché, formularios validados y estado del cliente, sin reinventar la rueda.",
        {"size": 16, "color": RGBColor(0xCB, 0xD5, 0xE1)})])
_text(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5),
      [("Programación / Diseño Web · Material de clase", {"size": 14, "color": MUTED})])
_text(s, Inches(0.9), Inches(6.95), Inches(11.5), Inches(0.4),
      [("Agosto 2026", {"size": 13, "color": MUTED})])

# ------------------------------------------------------------ slide 2 · agenda
_content_slide(
    "Agenda",
    subtitle="Qué vamos a ver y por qué",
    items=[
        ("bullet", [("El problema actual: ", {"bold": True}),
                    ("cómo conectamos el frontend a la API hoy (fetch manual, mock, estados repetidos).", {})]),
        ("bullet", [("Server State vs Client State: ", {"bold": True}),
                    ("dos tipos de estado que no se deben mezclar.", {})]),
        ("bullet", [("TanStack Query: ", {"bold": True}),
                    ("useQuery, useMutation, caché, invalidación y optimistic updates.", {})]),
        ("bullet", [("React Hook Form + Zod: ", {"bold": True}),
                    ("formularios gestionados con validación tipada en un solo esquema.", {})]),
        ("bullet", [("Zustand: ", {"bold": True}),
                    ("el estado del navegador (carrito, filtros) sin fricción.", {})]),
        ("bullet", [("Estados de UI: ", {"bold": True}),
                    ("loading, error y empty de una sola forma en todo el sistema.", {})]),
        ("bullet", [("Proyecto ejemplo ", {"bold": True}),
                    ("en Material-Clase/react-query-rhf-example y cómo aplicarlo a cada grupo.", {})]),
    ],
)

# ------------------------------------------------------------ slide 3 · problema
_content_slide(
    "El problema: cómo conectamos hoy",
    subtitle="Lo que se repite en los proyectos y por qué hay que cambiarlo",
    items=[
        ("bullet", [("fetch() suelto en cada componente ", {"bold": True}),
                    ("→ la lógica de datos queda dispersa y sin tipado.", {})]),
        ("bullet", [("Estados manuales con useEffect + useState ", {"bold": True}),
                    ("→ cada pantalla reinventa loading/error/vacío (y se olvida de alguno).", {})]),
        ("bullet", [("Fallback a mock data ", {"bold": True}),
                    ("→ la app funciona con datos falsos y la integración real queda sin probar.", {})]),
        ("bullet", [("Sin caché ", {"bold": True}),
                    ("→ cada navegación vuelve a pedir la lista entera al backend.", {})]),
        ("bullet", [("Validación duplicada ", {"bold": True}),
                    ("→ reglas en el HTML, en el submit y en el backend, sin una fuente única.", {})]),
        ("spacer", ""),
        ("note", "Estos tres problemas se resuelven con las tres librerías del taller."),
    ],
)

# ------------------------------------------------------------ slide 4 · server/client state
_content_slide(
    "Server State vs Client State",
    subtitle="El concepto base: hay dos tipos de estado y no se deben mezclar",
    items=[
        ("bullet", [("Server State (estado de servidor): ", {"bold": True, "color": ACCENT}),
                    ("los datos que viven en la API (productos, usuarios, vehículos…).", {})]),
        ("bullet", [("Es asíncrono, compartido, y puede quedar desactualizado ", {}),
                    ("si otro usuario cambió los datos.", {})]),
        ("bullet", [("Client State (estado de cliente): ", {"bold": True, "color": ACCENT}),
                    ("lo que solo existe en el navegador (carrito, filtro activo, tema, menú abierto).", {})]),
        ("bullet", [("Mezclarlos ", {"bold": True}),
                    ("es la causa de bugs clásicos: datos que no se actualizan o UI que no refleja la API.", {})]),
        ("spacer", ""),
        ("bullet", [("La división de trabajo: ", {"bold": True}),
                    ("TanStack Query → server state · Zustand → client state.", {})]),
    ],
)

# ------------------------------------------------------------ slide 5 · TanStack Query intro
_content_slide(
    "TanStack Query — qué es y por qué",
    subtitle="La capa de datos del frontend",
    items=[
        ("bullet", [("Librería estándar para datos asíncronos en React ", {"bold": True}),
                    ("(fetching, caché, sincronización y actualizaciones).", {})]),
        ("bullet", [("useQuery ", {"bold": True}),
                    ("para leer (GET) con caché y estados listos.", {})]),
        ("bullet", [("useMutation ", {"bold": True}),
                    ("para escribir (POST/PUT/DELETE) con loading, error y rollback.", {})]),
        ("bullet", [("Invalidación de caché: ", {"bold": True}),
                    ("tras una mutación, la lista se vuelve a pedir automáticamente.", {})]),
        ("bullet", [("Por qué: ", {"bold": True}),
                    ("deja de esconder la integración detrás de mocks y resuelve los estados de design.md en una línea.", {})]),
        ("spacer", ""),
        ("code", "const { data, isLoading, isError, refetch } = useQuery({\n"
                 "  queryKey: ['productos'],\n"
                 "  queryFn: () => apiRequest('/productos'),\n"
                 "});"),
    ],
)

# ------------------------------------------------------------ slide 6 · useQuery estados
_content_slide(
    "useQuery — los estados resueltos",
    subtitle="Cada pantalla de listado debería contemplar: carga, error y vacío",
    items=[
        ("bullet", [("isLoading: ", {"bold": True}),
                    ("primer pedido en curso → skeleton.", {})]),
        ("bullet", [("isError + error.message: ", {"bold": True}),
                    ("fallo de red o 404/500 → mensaje claro + botón reintentar.", {})]),
        ("bullet", [("data vacía: ", {"bold": True}),
                    ("0 resultados → empty state que invita a crear.", {})]),
        ("bullet", [("data con contenido: ", {"bold": True}),
                    ("la lista, servida desde la caché si ya se pidió antes.", {})]),
        ("spacer", ""),
        ("bullet", [("Por qué: ", {"bold": True}),
                    ("es el mismo patrón en todas las entidades (consistencia) y el código de estados no se duplica.", {})]),
        ("note", "En el proyecto ejemplo esto está en app/productos/page.tsx (LoadingSkeleton, ErrorState, EmptyState)."),
    ],
)

# ------------------------------------------------------------ slide 7 · mutations + invalidación
_content_slide(
    "useMutation + invalidación de caché",
    subtitle="Escribir en la API y mantener la UI sincronizada",
    items=[
        ("bullet", [("useMutation ", {"bold": True}),
                    ("maneja el ciclo de una escritura: isPending, isError, onSuccess.", {})]),
        ("bullet", [("Al crear/editar un producto: ", {"bold": True}),
                    ("onSuccess → invalidateQueries(['productos']).", {})]),
        ("bullet", [("invalidateQueries ", {"bold": True}),
                    ("marca la lista como vieja y la vuelve a pedir: la UI refleja la API.", {})]),
        ("spacer", ""),
        ("code", "const crear = useMutation({\n"
                 "  mutationFn: (datos) => apiRequest('/productos', { method: 'POST', body: JSON.stringify(datos) }),\n"
                 "  onSuccess: () => queryClient.invalidateQueries({ queryKey: ['productos'] }),\n"
                 "});"),
        ("spacer", ""),
        ("bullet", [("Por qué: ", {"bold": True}),
                    ("sin esto, crear un elemento no aparecería en el listado hasta recargar la página.", {})]),
    ],
)

# ------------------------------------------------------------ slide 8 · optimistic
_content_slide(
    "Actualización optimista (optimistic updates)",
    subtitle="UI instantánea con rollback si el servidor falla",
    items=[
        ("bullet", [("onMutate: ", {"bold": True}),
                    ("modificamos la caché al instante (el elemento desaparece de la UI ya).", {})]),
        ("bullet", [("onError: ", {"bold": True}),
                    ("si la API rechaza, restauramos el estado anterior (rollback).", {})]),
        ("bullet", [("onSettled: ", {"bold": True}),
                    ("al final siempre se reconcilia contra el servidor (invalidate).", {})]),
        ("spacer", ""),
        ("code", "const eliminar = useMutation({\n"
                 "  mutationFn: (id) => apiRequest(`/productos/${id}`, { method: 'DELETE' }),\n"
                 "  onMutate: (id) => { cancelarQueries(); quitarDeLaCache(id); return { anterior }; },\n"
                 "  onError: (_e, _id, ctx) => restaurar(ctx.anterior),\n"
                 "  onSettled: () => invalidar(),\n"
                 "});"),
        ("spacer", ""),
        ("bullet", [("Por qué: ", {"bold": True}),
                    ("la app se siente rápida y, si falla, vuelve sola al estado correcto.", {})]),
    ],
)

# ------------------------------------------------------------ slide 9 · RHF
_content_slide(
    "React Hook Form — por qué",
    subtitle="Formularios gestionados sin estado manual",
    items=[
        ("bullet", [("Maneja valores, errores por campo, touched y submit ", {"bold": True}),
                    ("sin useState por input.", {})]),
        ("bullet", [("Registro declarativo: ", {"bold": True}),
                    ("register('nombre') + valueAsNumber para los números.", {})]),
        ("bullet", [("Menos re-renders: ", {"bold": True}),
                    ("solo se actualiza lo que cambia.", {})]),
        ("bullet", [("Deshabilita el botón mientras se envía (isPending) ", {"bold": True}),
                    ("y cancela el submit si ya está en curso.", {})]),
        ("spacer", ""),
        ("bullet", [("Por qué: ", {"bold": True}),
                    ("los formularios son la parte con más bugs; gestionarlos de una forma consistente los vuelve predecibles.", {})]),
        ("note", "Componente ProductoForm.tsx: el mismo formulario sirve para alta y edición."),
    ],
)

# ------------------------------------------------------------ slide 10 · Zod
_content_slide(
    "Zod — validación tipada compartida",
    subtitle="Las reglas y los tipos en un solo lugar",
    items=[
        ("bullet", [("Esquema = fuente única de verdad: ", {"bold": True}),
                    ("reglas (min, max, positive, int) y tipos inferidos (z.infer).", {})]),
        ("bullet", [("z.number({ message }) ", {"bold": True}),
                    ("da mensajes en español por campo.", {})]),
        ("bullet", [("Tipa el submit: ", {"bold": True}),
                    ("los valores del formulario llegan tipados a la mutación.", {})]),
        ("spacer", ""),
        ("code", "export const productoFormSchema = z.object({\n"
                 "  nombre: z.string().min(3, 'El nombre debe tener al menos 3 caracteres'),\n"
                 "  precio: z.number({ message: 'Ingresá un precio válido' }).positive(),\n"
                 "  stock: z.number().int().nonnegative(),\n"
                 "});\n"
                 "export type ProductoFormValues = z.infer<typeof productoFormSchema>;"),
        ("spacer", ""),
        ("bullet", [("Por qué: ", {"bold": True}),
                    ("la validación deja de estar duplicada en HTML + JS + backend y el tipado acompaña al dato.", {})]),
    ],
)

# ------------------------------------------------------------ slide 11 · RHF + Zod
_content_slide(
    "React Hook Form + Zod (el puente)",
    subtitle="zodResolver conecta el esquema con el formulario",
    items=[
        ("bullet", [("El resolver valida en cada cambio y alimenta formState.errors ", {"bold": True}),
                    ("con los mensajes de Zod.", {})]),
        ("bullet", [("Los errores se muestran por campo ", {"bold": True}),
                    ("debajo de cada input, no en un bloque genérico.", {})]),
        ("spacer", ""),
        ("code", "const { register, handleSubmit, formState: { errors } } =\n"
                 "  useForm<ProductoFormValues>({\n"
                 "    resolver: zodResolver(productoFormSchema),\n"
                 "    defaultValues,  // precargado en edición\n"
                 "  });\n\n"
                 "<form onSubmit={handleSubmit(onSubmit)}>\n"
                 "  <input {...register('nombre')} />\n"
                 "  {errors.nombre && <p>{errors.nombre.message}</p>}\n"
                 "</form>"),
        ("spacer", ""),
        ("bullet", [("Por qué: ", {"bold": True}),
                    ("la misma validación tipada que define el contrato es la que corre en el navegador.", {})]),
    ],
)

# ------------------------------------------------------------ slide 12 · Zustand
_content_slide(
    "Zustand — estado del cliente",
    subtitle="Lo que es solo del navegador no debe ir a la API",
    items=[
        ("bullet", [("Store global mínimo con hook: ", {"bold": True}),
                    ("useCarrito → items, agregar, quitar, limpiar.", {})]),
        ("bullet", [("Sin Provider: ", {"bold": True}),
                    ("se importa el hook y listo.", {})]),
        ("bullet", [("Suscriptores selectivos: ", {"bold": True}),
                    ("useCarrito((s) => s.items.length) solo re-renderiza el contador.", {})]),
        ("spacer", ""),
        ("code", "export const useCarrito = create((set) => ({\n"
                 "  items: [],\n"
                 "  agregar: (id) => set((s) => ({ items: [...s.items, id] })),\n"
                 "  quitar: (id) => set((s) => ({ items: s.items.filter((i) => i !== id) })),\n"
                 "}));\n\n"
                 "const cantidad = useCarrito((s) => s.items.length);"),
        ("spacer", ""),
        ("bullet", [("Por qué: ", {"bold": True}),
                    ("separar client state de server state evita el problema del slide 4.", {})]),
    ],
)

# ------------------------------------------------------------ slide 13 · estados UI
_content_slide(
    "Estados de UI en la práctica",
    subtitle="Lo que pedía design.md: cada pantalla con carga, error, vacío y datos",
    items=[
        ("bullet", [("LoadingSkeleton ", {"bold": True}),
                    ("· ErrorState con reintentar ", {"bold": True}),
                    ("· EmptyState ", {"bold": True}),
                    ("· datos.", {})]),
        ("bullet", [("Se resuelven una vez en componentes de UI ", {"bold": True}),
                    ("(components/ui/Estado.tsx) y se reutilizan en todas las entidades.", {})]),
        ("bullet", [("Con useQuery ", {"bold": True}),
                    ("isLoading / isError / data.length === 0 salen gratis; solo resta mostrarlos.", {})]),
        ("spacer", ""),
        ("bullet", [("Por qué: ", {"bold": True}),
                    ("consistencia visual y menos decisiones por pantalla: el patrón se copia, no se reinventa.", {})]),
    ],
)

# ------------------------------------------------------------ slide 14 · arquitectura
_content_slide(
    "Estructura del proyecto ejemplo",
    subtitle="Material-Clase/react-query-rhf-example",
    items=[
        ("code", "backend/   → API mínima FastAPI con /api/productos (CRUD)\n"
                 "frontend/\n"
                 "  lib/api.ts            cliente HTTP tipado (un solo fetch)\n"
                 "  lib/providers.tsx     QueryClientProvider global\n"
                 "  lib/store.ts          Zustand: carrito (client state)\n"
                 "  schemas/producto.ts   esquema Zod + tipos inferidos\n"
                 "  queries/productos.ts  hooks useQuery / useMutation\n"
                 "  components/ProductoForm.tsx   RHF + Zod (alta y edición)\n"
                 "  app/productos/        listado · nuevo · [id]/editar"),
        ("spacer", ""),
        ("bullet", [("Separación por responsabilidad: ", {"bold": True}),
                    ("datos (queries), contratos (schemas), UI (components) y páginas (app).", {})]),
        ("bullet", [("El frontend consume /api mediante rewrites ", {"bold": True}),
                    ("de next.config → localhost:8000.", {})]),
    ],
)

# ------------------------------------------------------------ slide 15 · buenas prácticas
_content_slide(
    "Buenas prácticas y patrones",
    subtitle="Lo que se espera ver en los proyectos de los grupos",
    items=[
        ("bullet", [("Una capa de servicios/queries ", {"bold": True}),
                    ("y un cliente HTTP tipado; nunca fetch() suelto en componentes.", {})]),
        ("bullet", [("Esquemas Zod como contrato ", {"bold": True}),
                    ("compartido entre formularios y llamadas a la API.", {})]),
        ("bullet", [("Estados de UI centralizados ", {"bold": True}),
                    ("(Skeleton / ErrorState / EmptyState) reutilizables.", {})]),
        ("bullet", [("Invalidación después de cada mutación ", {"bold": True}),
                    ("y optimistic updates solo en acciones críticas (eliminar).", {})]),
        ("bullet", [("Mensajes de error en español por campo ", {"bold": True}),
                    ("y botón deshabilitado mientras se procesa.", {})]),
    ],
)

# ------------------------------------------------------------ slide 16 · aplicación a cada grupo
_content_slide(
    "Cómo se aplica a cada proyecto",
    subtitle="El tema cierra la brecha de los grupos",
    items=[
        ("bullet", [("G2 (Auto Market): ", {"bold": True}),
                    ("reemplazar MOCK_VEHICLES y filtros locales por useQuery conectado a /api/v1/vehicles.", {})]),
        ("bullet", [("G1 (Cafecito): ", {"bold": True}),
                    ("migrar los updates optimistas de mesas a useMutation con onMutate/onError.", {})]),
        ("bullet", [("G3 (Restaurante): ", {"bold": True}),
                    ("ya usa TanStack Query + RHF + Zod: sirve de referencia del patrón.", {})]),
        ("bullet", [("G4 (Agro): ", {"bold": True}),
                    ("arrancar el frontend con esta estructura desde cero.", {})]),
        ("spacer", ""),
        ("note", "El objetivo: los 4 grupos terminan con fetching con caché, formularios validados y estados de UI consistentes."),
    ],
)

# ------------------------------------------------------------ slide 17 · cierre
s = _slide()
_rect(s, 0, 0, SW, SH, DARK)
_rect(s, Inches(0.9), Inches(2.5), Inches(1.6), Inches(0.09), ACCENT)
_text(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(1.0),
      [("Para llevar", {"bold": True, "size": 36, "color": WHITE})])
_text(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(1.6),
      [("TanStack Query para los datos · React Hook Form + Zod para los formularios · Zustand para la UI.",
        {"size": 20, "color": ACCENT_SOFT})])
_text(s, Inches(0.9), Inches(4.5), Inches(11.5), Inches(1.2),
      [("Proyecto de práctica: Material-Clase/react-query-rhf-example\n"
        "Presentación: tanstack-query-rhf.pptx",
        {"size": 16, "color": RGBColor(0xCB, 0xD5, 0xE1)})])
_text(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5),
      [("¡A programar!", {"size": 14, "color": MUTED})])

# ---------------------------------------------------------------- guardar
OUT = "tanstack-query-rhf.pptx"
prs.save(OUT)
print(f"OK -> {OUT} ({len(prs.slides.__iter__.__self__._sldIdLst)} diapositivas)")